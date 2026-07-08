#!/usr/bin/env python3
"""Jungle_Notes PowerPoint presentation generator (13 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

W = Inches(10)
H = Inches(5.625)

BG = RGBColor(0x0F, 0x1B, 0x2D)
BG2 = RGBColor(0x1A, 0x2D, 0x4A)
CARD = RGBColor(0x16, 0x22, 0x40)
CARD2 = RGBColor(0x11, 0x1B, 0x33)

GOLD = RGBColor(0xD4, 0xA0, 0x17)
GOLD_B = RGBColor(0xF5, 0xC8, 0x42)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
BLUE_L = RGBColor(0x4A, 0x90, 0xD9)

WHT = RGBColor(0xFF, 0xFF, 0xFF)
MUT = RGBColor(0xB0, 0xBE, 0xC5)
MDM = RGBColor(0x78, 0x90, 0x9C)

S = RGBColor(0x16, 0x22, 0x40)  # shadow fill

LOGO_PATH = Path("public/images/utn-logo.jpg")
OUTPUT = "Jungle_Notes_Presentacion.pptx"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H




def bg(slide):
    """Add dark gradient background."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 0, 0, W, H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG
    # python-pptx doesn't support gradients, so solid fill
    sh.line.fill.background()


def box(slide, x, y, w, h, fill=CARD, outline=False):
    """Add a card shape."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if outline:
        sh.line.color.rgb = GOLD
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(slide, x, y, w, h, text, size=12, color=WHT, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb


def mul(slide, x, y, w, h, items, size=11, color=MUT, bullet_color=GOLD_B):
    """Add a bulleted text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u25c6  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tb


def dot(slide, cx, cy, d, color=GOLD):
    """Add a circle."""
    r = d // 2
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = S
    sh.line.color.rgb = color
    sh.line.width = Pt(2)
    return sh


def heading(slide, text):
    """Add slide title."""
    txt(slide, Inches(0.3), Inches(0.15), Inches(9), Inches(0.7), text, 24, WHT, True)


def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg(slide)
    return slide


def build_title():
    sl = new_slide()
    cx = W // 2

    if LOGO_PATH.exists():
        sl.shapes.add_picture(str(LOGO_PATH), cx - Inches(0.36), Inches(0.1), Inches(0.72), Inches(0.72))

    txt(sl, Inches(0.5), Inches(1.2), Inches(9), Inches(0.2), "UNIVERSIDAD TECNOLÓGICA NACIONAL · PROGRAMACIÓN 3", 10, BLUE_L, True, PP_ALIGN.CENTER)
    txt(sl, Inches(0.5), Inches(1.5), Inches(9), Inches(0.7), "Jungle Notes", 36, WHT, True, PP_ALIGN.CENTER)
    txt(sl, Inches(0.5), Inches(2.3), Inches(9), Inches(0.4), "Sistema Web de Gestión de Apuntes Académicos", 14, GOLD_B, False, PP_ALIGN.CENTER)

    cw = Inches(2.8)
    ch = Inches(1.4)
    mg = Inches(0.3)
    gy = Inches(3.4)
    data = [
        ("Institución", "Universidad Tecnológica\nNacional (UTN)\nSede Formosa, 2026"),
        ("Cátedra", "Programación 3\nProf. Lic. Enrique F. Aguirre,\nLic. Germán Raúl Rafael Crozy"),
        ("Equipo Iceberg", "Cardozo Emanuel\nCardozo Mauricio\nAntonelli Lucas\nBritez Santiago"),
    ]
    for i, (title, body) in enumerate(data):
        x = mg + i * (cw + Inches(0.15))
        box(sl, x, gy, cw, ch)
        txt(sl, x + Inches(0.08), gy + Inches(0.06), cw - Inches(0.16), Inches(0.25), title, 10, GOLD_B, True)
        txt(sl, x + Inches(0.08), gy + Inches(0.35), cw - Inches(0.16), ch - Inches(0.4), body, 9, MUT)


def build_compare(lt, lb, rt, rb, bn):
    sl = new_slide()
    heading(sl, "Problema vs Solución")

    mg = Inches(0.25)
    gw = int((W - 2 * mg - Inches(0.2)) / 2)
    gy = Inches(0.9)
    gh = Inches(3.5)

    for i, (title, color, bullets) in enumerate([(lt, BLUE, lb), (rt, GOLD, rb)]):
        x = mg + i * (gw + Inches(0.2))
        box(sl, x, gy, gw, gh)
        txt(sl, x + Inches(0.12), gy + Inches(0.12), gw - Inches(0.24), Inches(0.3), title, 18, WHT, True, PP_ALIGN.CENTER)
        dot(sl, x + gw // 2, gy + Inches(0.9), Inches(0.6), color)
        mul(sl, x + Inches(0.14), gy + Inches(1.6), gw - Inches(0.28), gh - Inches(1.7), bullets, 11, MUT, color)

    by = gy + gh + Inches(0.15)
    bh = Inches(0.6)
    box(sl, mg, by, W - 2 * mg, bh, CARD2, True)
    txt(sl, mg + Inches(0.2), by + Inches(0.06), W - 2 * mg - Inches(0.4), bh - Inches(0.12), bn, 12, WHT, True, PP_ALIGN.CENTER)


def build_grid(title, quads):
    sl = new_slide()
    heading(sl, title)

    mg = Inches(0.25)
    gap = Inches(0.15)
    ty = Inches(1.0)
    cw = int((W - 2 * mg - gap) / 2)
    ch = int((H - ty - Inches(0.25) - gap) / 2)

    for i, (c, r) in enumerate([(0, 0), (1, 0), (0, 1), (1, 1)]):
        q = quads[i]
        x = mg + c * (cw + gap)
        y = ty + r * (ch + gap)
        box(sl, x, y, cw, ch)
        ic = GOLD if (c + r) % 2 == 0 else BLUE
        dot(sl, x + Inches(0.35), y + ch // 2, Inches(0.5), ic)
        txt(sl, x + Inches(0.7), y + Inches(0.12), cw - Inches(0.85), Inches(0.3), q["title"], 13, WHT, True)
        txt(sl, x + Inches(0.7), y + Inches(0.5), cw - Inches(0.85), ch - Inches(0.6), q["desc"], 10, MDM)


def build_stack(title, layers):
    sl = new_slide()
    heading(sl, title)

    cx = W // 2
    dw = Inches(2.8)
    dh = Inches(0.55)
    y = Inches(1.0)
    colors = [GOLD, BLUE, GOLD]

    for i, ly in enumerate(layers):
        cl = colors[i % len(colors)]
        d = sl.shapes.add_shape(MSO_SHAPE.DIAMOND, cx - dw // 2, y, dw, dh)
        d.fill.solid()
        d.fill.fore_color.rgb = S
        d.line.color.rgb = cl
        d.line.width = Pt(2)
        txt(sl, cx - dw // 2 + Inches(0.12), y + dh // 2 - Inches(0.12), dw - Inches(0.24), Inches(0.24), ly["name"], 11, WHT, True, PP_ALIGN.CENTER)

        txty = y + dh + Inches(0.08)
        colw = Inches(2.8)
        colgap = Inches(0.2)
        colx = cx - colw - colgap // 2

        for j, col in enumerate([ly["col_a"], ly["col_b"]]):
            tx = colx + j * (colw + colgap)
            txt(sl, tx, txty, colw, Inches(0.2), col[0], 9, cl, True)
            txt(sl, tx, txty + Inches(0.2), colw, Inches(0.35), col[1], 8, MDM)

        y = txty + Inches(0.6)
        if i < len(layers) - 1:
            a = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, cx - Inches(0.08), y, Inches(0.16), Inches(0.14))
            a.fill.solid()
            a.fill.fore_color.rgb = BLUE
            a.line.fill.background()
            y += Inches(0.25)


def build_cycle(title, nodes, ct):
    sl = new_slide()
    heading(sl, title)

    cx = W // 2 + Inches(0.15)
    cy = Inches(3.0)
    rx = Inches(2.8)
    ry = Inches(1.2)
    nd = Inches(1.15)
    angles = {"l": 180, "t": 270, "r": 0, "b": 90}
    centers = {}
    for k, a in angles.items():
        import math
        rd = math.radians(a)
        centers[k] = (int(cx + rx * math.cos(rd)), int(cy + ry * math.sin(rd)))

    order = ["l", "t", "r", "b"]
    for i in range(4):
        a = order[i]
        b = order[(i + 1) % 4]
        connector = sl.shapes.add_connector(1, centers[a][0], centers[a][1], centers[b][0], centers[b][1])  # type 1 = straight
        connector.line.color.rgb = GOLD
        connector.line.width = Pt(2)

    for i, k in enumerate(order):
        ndata = nodes[i]
        px, py = centers[k]
        c = sl.shapes.add_shape(MSO_SHAPE.OVAL, px - nd // 2, py - nd // 2, nd, nd)
        c.fill.solid()
        c.fill.fore_color.rgb = CARD
        c.line.color.rgb = BLUE
        c.line.width = Pt(2)
        txt(sl, px - nd // 2 + Inches(0.08), py - Inches(0.25), nd - Inches(0.16), Inches(0.3), f"{ndata['number']}. {ndata['label']}", 10, WHT, True, PP_ALIGN.CENTER)
        if ndata.get("tagline"):
            txt(sl, px - nd // 2 + Inches(0.08), py + Inches(0.06), nd - Inches(0.16), Inches(0.35), ndata["tagline"], 7, MDM, False, PP_ALIGN.CENTER)

    cw = Inches(2.2)
    ch = Inches(0.65)
    box(sl, cx - cw // 2, cy - ch // 2, cw, ch, CARD2)
    txt(sl, cx - cw // 2 + Inches(0.08), cy - ch // 2 + Inches(0.08), cw - Inches(0.16), ch - Inches(0.16), ct, 9, MUT, False, PP_ALIGN.CENTER)


def build_cards_row(title, cards, hi=None):
    sl = new_slide()
    heading(sl, title)

    n = len(cards)
    mg = Inches(0.25)
    gap = Inches(0.12)
    cw = int((W - 2 * mg - gap * (n - 1)) / n)
    y = Inches(1.0)
    ch = Inches(4.2)

    for i, card in enumerate(cards):
        x = mg + i * (cw + gap)
        hl = i == hi
        box(sl, x, y, cw, ch, CARD2 if hl else CARD, hl)
        ic = GOLD_B if hl else BLUE
        dot(sl, x + cw // 2, y + Inches(0.6), Inches(0.65), ic)
        txt(sl, x + Inches(0.08), y + Inches(1.1), cw - Inches(0.16), Inches(0.35), card["title"], 13, WHT, True, PP_ALIGN.CENTER)
        txt(sl, x + Inches(0.12), y + Inches(1.5), cw - Inches(0.24), ch - Inches(1.7), card["desc"], 9, MDM, False, PP_ALIGN.CENTER)


def build_feature(title, ft, items):
    sl = new_slide()
    heading(sl, title)

    mg = Inches(0.25)
    cw = W - 2 * mg
    cy = Inches(0.9)
    ch = Inches(4.3)

    box(sl, mg, cy, cw, ch)
    txt(sl, mg + Inches(0.14), cy + Inches(0.12), cw - Inches(0.28), Inches(0.3), ft, 18, GOLD_B, True)

    body = ""
    for it in items:
        body += f"\u25c6  {it['title']}\n{it['desc']}\n\n"
    txt(sl, mg + Inches(0.14), cy + Inches(0.55), cw - Inches(0.28), ch - Inches(0.7), body.strip(), 10, MUT)


def build_simple_text(title, body):
    sl = new_slide()
    heading(sl, title)

    mg = Inches(0.25)
    cw = W - 2 * mg
    cy = Inches(0.9)
    ch = Inches(4.3)

    box(sl, mg, cy, cw, ch)
    txt(sl, mg + Inches(0.14), cy + Inches(0.12), cw - Inches(0.28), ch - Inches(0.24), body, 11, MUT)


def build_close():
    sl = new_slide()
    txt(sl, Inches(0.5), Inches(1.4), Inches(9), Inches(0.7), "\u00a1Gracias!", 42, WHT, True, PP_ALIGN.CENTER)
    txt(sl, Inches(0.5), Inches(2.2), Inches(9), Inches(0.4), "\u00bfPreguntas?", 18, GOLD_B, False, PP_ALIGN.CENTER)
    txt(sl, Inches(0.5), Inches(3.0), Inches(9), Inches(0.3), "Equipo Iceberg \u00b7 Programaci\u00f3n III \u00b7 UTN Formosa", 12, MDM, False, PP_ALIGN.CENTER)
    dot(sl, W // 2, Inches(4.15), Inches(0.6), GOLD)
    txt(sl, W // 2 - Inches(0.3), Inches(4.0), Inches(0.6), Inches(0.3), "?", 24, WHT, True, PP_ALIGN.CENTER)


# -- Build deck --
build_title()

build_compare(
    "El Desaf\u00edo",
    ["Apuntes dispersos en archivos sueltos.", "Sin organizaci\u00f3n por materia.", "Dificultad para encontrar material de estudio."],
    "La Soluci\u00f3n",
    ["Sistema centralizado con b\u00fasqueda inteligente.", "Organizaci\u00f3n autom\u00e1tica por materia.", "Acceso r\u00e1pido desde cualquier dispositivo."],
    "Jungle Notes centraliza, organiza y facilita la b\u00fasqueda de apuntes acad\u00e9micos."
)

build_grid("Objetivos del Proyecto", [
    {"title": "Gesti\u00f3n de Apuntes", "desc": "CRUD completo para subir, buscar y\neliminar archivos PDF/Word\norganizados por materia."},
    {"title": "Dashboard Principal", "desc": "Panel con m\u00e9tricas globales de\napuntes y materias, m\u00e1s actividad\nreciente."},
    {"title": "Bloc de Notas", "desc": "Editor r\u00e1pido de notas personales\nintegrado en el sistema."},
    {"title": "Materias y Comentarios", "desc": "Administraci\u00f3n de materias con\ndetalle de apuntes y sistema de\ncomentarios."},
])

build_feature("Gesti\u00f3n de Materias", "Materias", [
    {"title": "CRUD Completo", "desc": "Crear, editar, ver y eliminar materias con nombre, descripci\u00f3n y a\u00f1o."},
    {"title": "Apuntes Asociados", "desc": "Cada materia muestra sus apuntes relacionados con b\u00fasqueda incluida."},
    {"title": "Almacenamiento Estructurado", "desc": "Archivos organizados en storage/app/public/apuntes/materia_{id}/."},
])

build_cards_row("Calendario Acad\u00e9mico", [
    {"title": "FullCalendar", "desc": "Calendario interactivo con vista mensual, semanal y diaria."},
    {"title": "Eventos", "desc": "Crear, editar y eliminar ex\u00e1menes, presentaciones y eventos acad\u00e9micos."},
    {"title": "Feriados", "desc": "Feriados nacionales mostrados como fondo en el calendario."},
])

build_simple_text("Bloc de Notas", "\u25c6  Notas R\u00e1pidas\n   Editor simple para tomar notas personales durante el estudio.\n\n\u25c6  Persistencia Autom\u00e1tica\n   Las notas se guardan autom\u00e1ticamente al escribir.\n\n\u25c6  Integraci\u00f3n\n   Acceso directo desde el men\u00fa principal del sistema.")

build_simple_text("Sistema de Comentarios", "\u25c6  Comentarios en Tiempo Real\n   Los comentarios se actualizan sin recargar la p\u00e1gina.\n\n\u25c6  Por Apunte\n   Cada apunte tiene su propio hilo de comentarios.\n\n\u25c6  Colaborativo\n   Facilita la discusi\u00f3n y el intercambio de ideas entre compa\u00f1eros.")

build_stack("Stack Tecnol\u00f3gico", [
    {"name": "Frontend", "col_a": ["Tailwind CSS v4", "Framework de utilidades para dise\u00f1o moderno y responsive."], "col_b": ["Livewire v4 + Flux UI v2", "Componentes din\u00e1micos con Alpine.js integrado."]},
    {"name": "Backend", "col_a": ["Laravel 12", "Framework para rutas, validaciones y seguridad."], "col_b": ["PHP 8.3", "Lenguaje principal del lado del servidor."]},
    {"name": "Base de Datos", "col_a": ["SQLite / MySQL", "Almacenamiento relacional consistente."], "col_b": ["Eloquent ORM", "Consultas expresivas y seguras con modelo MVC."]},
])

build_cycle("Arquitectura MVC en Livewire", [
    {"number": 1, "label": "El Usuario", "tagline": "Interact\u00faa v\u00eda navegador"},
    {"number": 2, "label": "Componente", "tagline": "Coordina l\u00f3gica y estado"},
    {"number": 3, "label": "Modelo", "tagline": "Gestiona los datos"},
    {"number": 4, "label": "Vista Blade", "tagline": "Renderiza el resultado"},
], "Livewire elimina la separaci\u00f3n tradicional Controller/Vista.")

build_cards_row("Entidades de la Base de Datos", [
    {"title": "Materia", "desc": "Agrupa apuntes por asignatura. Campos: nombre, descripci\u00f3n, a\u00f1o."},
    {"title": "Apunte", "desc": "Archivo PDF/Word con t\u00edtulo, descripci\u00f3n y materia asociada."},
    {"title": "Comentario", "desc": "Mensaje asociado a un apunte, con usuario y contenido."},
    {"title": "Evento", "desc": "Examen o presentaci\u00f3n con fecha, tipo y color en el calendario."},
    {"title": "Nota", "desc": "Nota personal con t\u00edtulo y contenido."},
], 1)

build_simple_text("Demo del Sistema", "\u25c6  Pantalla Principal\n   Dashboard con m\u00e9tricas y actividad reciente.\n\n\u25c6  Gesti\u00f3n de Apuntes\n   Subir, buscar y eliminar archivos con filtro por materia.\n\n\u25c6  Calendario Acad\u00e9mico\n   Visualizar y gestionar eventos acad\u00e9micos.\n\n\u25c6  Bloc de Notas\n   Tomar notas r\u00e1pidas durante el estudio.")

build_simple_text("Conclusiones", "\u25c6  Centralizaci\u00f3n\n   Un solo lugar para todos los apuntes, materias y eventos acad\u00e9micos.\n\n\u25c6  Tecnolog\u00eda Moderna\n   Laravel + Livewire + Tailwind: stack eficiente y productivo.\n\n\u25c6  Trabajo en Equipo\n   Desarrollo colaborativo con Git, issues y pull requests.\n\n\u25c6  Pr\u00f3ximos Pasos\n   Mejoras de seguridad, optimizaci\u00f3n de cach\u00e9 y m\u00e1s features.")

build_close()

prs.save(OUTPUT)
print(f"Presentación generada: {OUTPUT}")
